from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_box(shapes, left, top, width, height, color, text, font_size=Pt(12)):
    """輔助函數：建立帶有背景色的文字方塊"""
    shape = shapes.add_shape(1, left, top, width, height)  # 1 是長方形
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(200, 200, 200)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = RGBColor(0, 0, 0)
    return shape


# 建立簡報物件
prs = Presentation()

# 設定投影片寬高 (16:9)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # 使用空白佈局
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

# 1. 頂部深綠色標題列
header = shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = RGBColor(30, 120, 90)
header.line.fill.background()

title_shape = shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
tf = title_shape.text_frame
p = tf.paragraphs[0]
p.text = "產業升級創新平台輔導計畫 進度說明"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(255, 255, 255)

# 2. 主標題與效益區
shapes.add_textbox(Inches(0.5), Inches(1), Inches(10), Inches(0.4)).text_frame.paragraphs[0].text = "國產高效能氮化鎵電動多功能農機平台"
benefits_text = "效益：透過晶創農機計畫落地輔導，帶動國內指標廠商投入，建立國內第一條 GaN 馬達驅動器產線...預期產值 3,250 ~ 5,250 萬元。"
add_box(shapes, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1), RGBColor(255, 245, 180), benefits_text, Pt(11))

# 3. 中間三個資訊方塊
# 左：核心推進
add_box(shapes, Inches(0.5), Inches(2.7), Inches(3.8), Inches(3.5), RGBColor(240, 250, 245),
        "核心推進重點\n\n● 完成計畫書初稿\n● 經費效益對齊\n● 聚焦技術定位", Pt(11))

# 中：驅動器
add_box(shapes, Inches(4.7), Inches(2.7), Inches(3.8), Inches(3.5), RGBColor(240, 250, 245),
        "國產高效能 GaN 馬達驅動器\n\n● 採用 GaN 功率元件\n● 支援 FOC 控制\n● 延伸至 AGV/AMR 市場", Pt(11))

# 右：農機平台
add_box(shapes, Inches(8.9), Inches(2.7), Inches(3.8), Inches(3.5), RGBColor(240, 250, 245),
        "國產高效能電動多功能農機平台\n\n● 高效率 GaN 驅動\n● 模組化平台設計\n● 具場域驗證基礎", Pt(11))

# 4. 底部進度區
bottom_box = shapes.add_shape(1, Inches(0.5), Inches(6.4), Inches(12.3), Inches(1))
bottom_box.fill.background()
bottom_box.line.color.rgb = RGBColor(150, 150, 150)
tf = bottom_box.text_frame
p1 = tf.paragraphs[0]
p1.text = "3月 推進完成重點：完成初稿、經費對齊、PMC 說明"
p2 = tf.add_paragraph()
p2.text = "4月 預期後續：補齊附件、最終校核、依時程送件"

# 存檔
prs.save('Project_Progress_Report.pptx')
print("PPT 已成功生成！")